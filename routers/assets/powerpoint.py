import os
from modules.logger_tool import initialise_logger
logger = initialise_logger(log_name="powerpoint", log_level=os.getenv("LOG_LEVEL"), log_dir=os.getenv("LOG_PATH"), log_format="default", runtime=True)

from fastapi import APIRouter, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
from pathlib import Path
import tempfile
from pptx import Presentation
from PIL import Image
import io
import base64
import traceback
import sys
import subprocess
from concurrent.futures import ThreadPoolExecutor, as_completed, TimeoutError
import asyncio
import psutil
import math
import time

router = APIRouter()

# Global semaphore to control total concurrent PowerPoint processing
MAX_CONCURRENT_PROCESSING = 4  # Adjust based on server capacity
processing_semaphore = asyncio.Semaphore(MAX_CONCURRENT_PROCESSING)

def calculate_optimal_workers():
    """Calculate optimal number of worker threads based on system resources."""
    cpu_count = os.cpu_count() or 4
    available_memory = psutil.virtual_memory().available
    memory_per_worker = 500 * 1024 * 1024  # 500MB per worker estimate
    
    # Calculate workers based on CPU and memory constraints
    cpu_based_workers = max(1, cpu_count - 1)  # Leave one core free
    memory_based_workers = max(1, int(available_memory / memory_per_worker))
    
    # Take the minimum of CPU and memory-based calculations
    optimal_workers = min(cpu_based_workers, memory_based_workers)
    
    # Cap at a reasonable maximum
    final_workers = min(optimal_workers, 8)  # Maximum 8 workers per process

    # Log resource information
    logger.info("Resource utilization:", {
        "total_cpus": cpu_count,
        "available_memory_gb": available_memory / (1024**3),
        "cpu_based_workers": cpu_based_workers,
        "memory_based_workers": memory_based_workers,
        "final_workers": final_workers
    })
    
    return final_workers

def extract_text_from_shape(shape):
    """Extract text from a PowerPoint shape."""
    if hasattr(shape, 'text') and shape.text.strip():
        return shape.text.strip()

    # Handle tables
    if shape.has_table:
        table_text = []
        for row in shape.table.rows:
            row_text = []
            row_text.extend(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                table_text.append('| ' + ' | '.join(row_text) + ' |')
        if table_text:
            # Add markdown table header separator
            table_text.insert(1, '|' + '---|' * (len(table_text[0].split('|')) - 2))
            return '\n'.join(table_text)

    # Handle grouped shapes
    if hasattr(shape, 'shapes'):
        group_text = []
        for subshape in shape.shapes:
            if text := extract_text_from_shape(subshape):
                group_text.append(text)
        return '\n'.join(group_text) if group_text else ''

    return ''

def extract_slide_text(slide):
    """Extract text from a PowerPoint slide and format as markdown."""
    slide_text = []

    # Extract title if present
    if slide.shapes.title and slide.shapes.title.text.strip():
        slide_text.append(f"# {slide.shapes.title.text.strip()}")

    # Process all shapes
    for shape in slide.shapes:
        if shape != slide.shapes.title:  # Skip title as we've already processed it
            if text := extract_text_from_shape(shape):
                slide_text.append(text)

    return '\n\n'.join(slide_text)

def process_slide(temp_dir: str, pdf_path: str, pptx_path: str, slide_info: tuple, timeout: int = 30) -> dict:
    """
    Worker function to process a single slide and enforce 16:9 aspect ratio.
    Args:
        temp_dir: Path to temporary directory
        pdf_path: Path to PDF file
        pptx_path: Path to PowerPoint file
        slide_info: Tuple of (index, slide_number)
        timeout: Maximum time in seconds to process a single slide
    Returns:
        dict: Processed slide information
    """
    i, slide_idx = slide_info
    slide_num = slide_idx + 1  # PDF pages are 1-indexed
    output_prefix = str(Path(temp_dir) / f"slide_{slide_num}")

    try:
        # Extract text from PowerPoint slide
        prs = Presentation(pptx_path)
        slide_text = extract_slide_text(prs.slides[slide_idx])

        # Convert PDF page to PNG with timeout
        process = subprocess.Popen(
            [
                'pdftoppm',
                '-png',
                '-singlefile',
                '-f',
                str(slide_num),
                '-l',
                str(slide_num),
                '-r',
                '600',  # High resolution for better quality
                pdf_path,
                output_prefix,
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE
        )
        
        try:
            stdout, stderr = process.communicate(timeout=timeout)
        except subprocess.TimeoutExpired:
            process.kill()
            raise TimeoutError(f"Slide {slide_num} processing timed out after {timeout} seconds")

        if process.returncode != 0:
            raise Exception(f"pdftoppm failed for slide {slide_num}: {stderr.decode()}")

        output_file = f"{output_prefix}.png"
        if not Path(output_file).exists():
            raise Exception(f"Could not find output file for slide {slide_num}")

        # Open and process the image
        with Image.open(output_file) as img:
            result = _process_image(img, i)
            if result['success']:
                result['meta'] = {
                    'text': slide_text,
                    'format': 'markdown'
                }
            return result
    except Exception as e:
        logger.error(f"Error processing slide {slide_num}: {str(e)}")
        return {
            "index": i,
            "error": str(e),
            "success": False,
        }

def _process_image(img: Image.Image, index: int) -> dict:
    """Process a single image, enforcing aspect ratio and size constraints."""
    try:
        # Enforce 16:9 aspect ratio
        target_aspect_ratio = 16 / 9
        img_aspect_ratio = img.width / img.height

        if img_aspect_ratio > target_aspect_ratio:  # Wider than 16:9
            new_width = int(img.height * target_aspect_ratio)
            offset = (img.width - new_width) // 2
            img = img.crop((offset, 0, offset + new_width, img.height))
        elif img_aspect_ratio < target_aspect_ratio:  # Taller than 16:9
            new_height = int(img.width / target_aspect_ratio)
            offset = (img.height - new_height) // 2
            img = img.crop((0, offset, img.width, offset + new_height))

        # Resize to target resolution (2560x1440)
        img = img.resize((2560, 1440), Image.Resampling.LANCZOS)

        # Convert to base64
        buffered = io.BytesIO()
        img.save(buffered, format="PNG", optimize=True)
        img_str = base64.b64encode(buffered.getvalue()).decode()

        return {
            "index": index,
            "data": f"data:image/png;base64,{img_str}",
            "success": True,
        }
    except Exception as e:
        logger.error(f"Error processing image for slide {index}: {str(e)}")
        return {
            "index": index,
            "error": str(e),
            "success": False,
        }

async def process_slides_in_chunks(temp_dir: str, pdf_path: str, pptx_path: str, visible_slides: list, chunk_size: int = 5):
    """Process slides in chunks to manage memory better."""
    all_processed_slides = []
    num_workers = calculate_optimal_workers()
    total_chunks = math.ceil(len(visible_slides) / chunk_size)
    
    logger.info("Starting slide processing:", {
        "total_slides": len(visible_slides),
        "chunk_size": chunk_size,
        "total_chunks": total_chunks,
        "workers_per_chunk": num_workers
    })

    # Process slides in chunks
    for chunk_index in range(0, len(visible_slides), chunk_size):
        chunk = visible_slides[chunk_index:chunk_index + chunk_size]
        processed_chunk = []
        current_chunk_num = (chunk_index // chunk_size) + 1

        logger.info(f"Processing chunk {current_chunk_num}/{total_chunks}", {
            "chunk_size": len(chunk),
            "chunk_start_index": chunk_index,
            "memory_usage_gb": psutil.Process().memory_info().rss / (1024**3)
        })

        start_time = time.time()
        with ThreadPoolExecutor(max_workers=num_workers) as executor:
            # Submit chunk of tasks
            future_to_slide = {
                executor.submit(
                    process_slide, temp_dir, pdf_path, pptx_path, slide_info
                ): slide_info
                for slide_info in chunk
            }

            # Process completed tasks as they finish
            for future in as_completed(future_to_slide):
                try:
                    result = future.result(timeout=60)  # Increased timeout to 60 seconds per slide
                    if result.get('success', False):
                        processed_chunk.append(result)
                    slide_info = future_to_slide[future]
                    logger.debug(f"Processed slide {slide_info[1] + 1}", {
                        "success": result.get('success', False),
                        "processing_time": time.time() - start_time
                    })
                except TimeoutError:
                    slide_info = future_to_slide[future]
                    logger.error(f"Timeout processing slide {slide_info[1] + 1}")
                except Exception as e:
                    slide_info = future_to_slide[future]
                    logger.error(f"Error processing slide {slide_info[1] + 1}: {str(e)}")

        chunk_time = time.time() - start_time
        logger.info(f"Completed chunk {current_chunk_num}/{total_chunks}", {
            "processed_slides": len(processed_chunk),
            "chunk_processing_time": chunk_time,
            "avg_time_per_slide": chunk_time / len(chunk) if chunk else 0
        })

        all_processed_slides.extend(processed_chunk)

        # Small delay between chunks to allow other tasks to process
        await asyncio.sleep(0.1)

    return all_processed_slides

@router.post("/convert")
async def convert_pptx_to_images(file: UploadFile = File(...)):
    try:
        async with processing_semaphore:  # Control concurrent processing
            start_time = time.time()
            # Log request details
            logger.info(
                "Received file upload request",
                {
                    "filename": file.filename,
                    "content_type": file.content_type,
                    "current_memory_usage_gb": psutil.Process()
                    .memory_info()
                    .rss
                    / (1024**3),
                    "cpu_percent": psutil.cpu_percent(interval=1),
                },
            )

            # Validate file
            if not file.filename.endswith('.pptx'):
                logger.error("Invalid file type")
                return JSONResponse({
                    "status": "error",
                    "message": "Invalid file type. Please upload a .pptx file"
                }, status_code=400)

            # Create a temporary directory to store the PowerPoint file
            with tempfile.TemporaryDirectory() as temp_dir:
                pptx_path = Path(temp_dir) / "presentation.pptx"
                logger.debug(f"Saving file to temporary path: {pptx_path}")

                try:
                    # Save uploaded file
                    content = await file.read()
                    logger.debug(f"Read file content, size: {len(content)} bytes")

                    with open(pptx_path, "wb") as buffer:
                        buffer.write(content)
                    logger.debug("File saved successfully")

                    if not pptx_path.exists() or pptx_path.stat().st_size == 0:
                        raise Exception("Failed to save file or file is empty")

                    # Open the presentation and get visible slides
                    prs = Presentation(str(pptx_path))
                    visible_slides = [
                        (i, slide_idx)
                        for i, (slide_idx, _) in enumerate(
                            (i, slide)
                            for i, slide in enumerate(prs.slides)
                            if not hasattr(slide, 'show') or slide.show
                        )
                    ]
                    num_slides = len(visible_slides)

                    if num_slides == 0:
                        logger.warning("No visible slides found in presentation")
                        return JSONResponse({
                            "status": "error",
                            "message": "No visible slides found in presentation"
                        }, status_code=400)

                    logger.info(f"Processing {num_slides} visible slides")

                    # Convert PowerPoint to PDF
                    pdf_path = Path(temp_dir) / "presentation.pdf"
                    logger.debug("Converting PowerPoint to PDF")

                    result = subprocess.run([
                        'soffice',
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', str(temp_dir),
                        str(pptx_path)
                    ], check=True, capture_output=True, text=True)

                    if not pdf_path.exists():
                        raise Exception("PDF file was not created")

                    logger.debug(f"PDF created successfully at {pdf_path}, size: {pdf_path.stat().st_size} bytes")

                    # Calculate chunk size based on number of slides
                    chunk_size = min(5, max(2, math.ceil(num_slides / 4)))
                    processed_slides = await process_slides_in_chunks(str(temp_dir), str(pdf_path), str(pptx_path), visible_slides, chunk_size)

                    if not processed_slides:
                        raise Exception("Failed to process any slides successfully")

                    # Sort slides by index
                    processed_slides.sort(key=lambda x: x['index'])

                    logger.info(f"Successfully processed {len(processed_slides)} slides")

                    # After processing all slides
                    total_time = time.time() - start_time
                    logger.info("PowerPoint processing completed", {
                        "total_processing_time": total_time,
                        "slides_processed": len(processed_slides),
                        "avg_time_per_slide": total_time / len(processed_slides) if processed_slides else 0,
                        "final_memory_usage_gb": psutil.Process().memory_info().rss / (1024**3)
                    })

                    return JSONResponse({
                        "status": "success",
                        "slides": processed_slides,
                        "processing_stats": {
                            "total_time": total_time,
                            "slides_processed": len(processed_slides),
                            "avg_time_per_slide": total_time / len(processed_slides) if processed_slides else 0
                        }
                    })

                except Exception as inner_error:
                    logger.error(f"Inner error: {str(inner_error)}")
                    logger.error(traceback.format_exc())
                    raise

    except Exception as e:
        logger.error(f"Error processing PowerPoint: {str(e)}")
        logger.error(f"Python version: {sys.version}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return JSONResponse({
            "status": "error",
            "message": f"Failed to process PowerPoint: {str(e)}"
        }, status_code=500)